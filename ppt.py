# import streamlit as st
# import tempfile
# import os
# import warnings
# from pptx import Presentation
# from pptx.util import Inches, Pt
# from pptx.enum.text import PP_ALIGN
# from PIL import Image
# from reportlab.lib.pagesizes import A4
# from reportlab.pdfgen import canvas
# import shutil

# # ======================
# # Configuration
# # ======================
# SLIDE_WIDTH = Inches(13.33)
# SLIDE_HEIGHT = Inches(7.5)
# PDF_PAGE_SIZE = A4
# TITLE_HEIGHT = Inches(0.8)
# GAP_BETWEEN = Inches(0.3)
# MARGIN = Inches(0.4)

# # ======================
# # Core Functions
# # ======================
# def scale_image(img_path, max_width, max_height):
#     with Image.open(img_path) as img:
#         width, height = img.size
#         aspect = height / width
        
#         width_based = min(max_width, (max_height - GAP_BETWEEN) / aspect)
#         height_based = min(max_height - GAP_BETWEEN, max_width * aspect)
        
#         if (width_based * (width_based * aspect)) > (height_based * (height_based / aspect)):
#             return width_based, width_based * aspect
#         else:
#             return height_based / aspect, height_based

# def create_ppt(photo_paths, output_path):
#     prs = Presentation()
#     prs.slide_width = SLIDE_WIDTH
#     prs.slide_height = SLIDE_HEIGHT
    
#     for img_path in photo_paths:
#         slide = prs.slides.add_slide(prs.slide_layouts[6])
        
#         title = os.path.splitext(os.path.basename(img_path))[0]
#         textbox = slide.shapes.add_textbox(
#             left=Inches(0.5),
#             top=Inches(0.3),
#             width=SLIDE_WIDTH - Inches(1),
#             height=TITLE_HEIGHT
#         )
#         p = textbox.text_frame.add_paragraph()
#         p.text = title
#         p.alignment = PP_ALIGN.CENTER
#         p.font.name = "Arial"
#         p.font.size = Pt(24)
#         p.font.bold = True
        
#         max_width = SLIDE_WIDTH - (2 * MARGIN)
#         max_height = SLIDE_HEIGHT - TITLE_HEIGHT - GAP_BETWEEN
        
#         img_width, img_height = scale_image(img_path, max_width, max_height)
        
#         left = (SLIDE_WIDTH - img_width) / 2
#         top = TITLE_HEIGHT + GAP_BETWEEN
        
#         slide.shapes.add_picture(
#             img_path,
#             left, top,
#             width=img_width,
#             height=img_height
#         )
    
#     prs.save(output_path)

# def create_pdf(photo_paths, output_path):
#     c = canvas.Canvas(output_path, pagesize=PDF_PAGE_SIZE)
#     page_width, page_height = PDF_PAGE_SIZE
#     title_font_size = 24
#     gap_points = GAP_BETWEEN.inches * 72
    
#     for img_path in photo_paths:
#         c.showPage()
        
#         title = os.path.splitext(os.path.basename(img_path))[0]
#         c.setFont("Helvetica-Bold", title_font_size)
#         title_y = page_height - 50
#         c.drawCentredString(page_width/2, title_y, title)
        
#         max_width = page_width - (MARGIN.inches * 72 * 2)
#         max_height = page_height - 70 - gap_points
        
#         with Image.open(img_path) as img:
#             img_width, img_height = scale_image(img_path, max_width, max_height)
        
#         x = (page_width - img_width) / 2
#         y = title_y - gap_points - img_height
        
#         c.drawImage(img_path, x, y, 
#                    width=img_width, 
#                    height=img_height,
#                    preserveAspectRatio=True)
    
#     c.save()

# # ======================
# # Streamlit UI
# # ======================
# st.set_page_config(page_title="Professional Photo Converter", layout="wide")
# st.title("📸 Smart Photo Converter")
# st.write("Upload photos to create perfect presentations!")

# uploaded_files = st.file_uploader(
#     "Drag & drop photos (multiple selection)",
#     type=["png", "jpg", "jpeg"],
#     accept_multiple_files=True
# )

# if uploaded_files:
#     temp_dir = tempfile.mkdtemp()
#     photo_paths = []
    
#     for uploaded_file in uploaded_files:
#         file_path = os.path.join(temp_dir, uploaded_file.name)
#         with open(file_path, "wb") as f:
#             f.write(uploaded_file.getbuffer())
#         photo_paths.append(file_path)
    
#     st.subheader("Photo Preview")
#     cols = st.columns(3)
#     for idx, path in enumerate(photo_paths):
#         try:
#             with warnings.catch_warnings():
#                 warnings.filterwarnings("ignore", category=DeprecationWarning)
#                 cols[idx % 3].image(
#                     path, 
#                     use_column_width=True,
#                     caption=os.path.basename(path)
#                 )
#         except Exception as e:
#             st.error(f"Error previewing {os.path.basename(path)}: {str(e)}")
    
#     st.subheader("Output Settings")
#     output_name = st.text_input("Base filename:", "MyPresentation")
#     output_format = st.radio("Format:", ["PPTX", "PDF", "Both"])
    
#     if st.button("🚀 Generate Files"):
#         with st.spinner("Creating outputs..."):
#             try:
#                 generated_files = []
                
#                 if output_format in ["PPTX", "Both"]:
#                     ppt_path = os.path.join(temp_dir, f"{output_name}.pptx")
#                     create_ppt(photo_paths, ppt_path)
#                     generated_files.append(ppt_path)
                
#                 if output_format in ["PDF", "Both"]:
#                     pdf_path = os.path.join(temp_dir, f"{output_name}.pdf")
#                     create_pdf(photo_paths, pdf_path)
#                     generated_files.append(pdf_path)
                
#                 if generated_files:
#                     st.success("✅ Files ready!")
#                     for file_path in generated_files:
#                         with open(file_path, "rb") as f:
#                             st.download_button(
#                                 label=f"📥 {os.path.basename(file_path)}",
#                                 data=f.read(),
#                                 file_name=os.path.basename(file_path),
#                                 mime="application/octet-stream" if "pptx" in file_path 
#                                      else "application/pdf"
#                             )
                
#             except Exception as e:
#                 st.error(f"⚠️ Error: {str(e)}")
#             finally:
#                 shutil.rmtree(temp_dir, ignore_errors=True)

# # requirements.txt remains same














# import streamlit as st
# import tempfile
# import os
# import warnings
# from pptx import Presentation
# from pptx.util import Inches, Pt
# from pptx.enum.text import PP_ALIGN
# from PIL import Image
# from reportlab.lib.pagesizes import A4
# from reportlab.pdfgen import canvas
# import shutil
# import csv
# import re

# # ======================
# # Configuration
# # ======================
# SLIDE_WIDTH = Inches(13.33)
# SLIDE_HEIGHT = Inches(7.5)
# PDF_PAGE_SIZE = A4
# TITLE_HEIGHT = Inches(0.8)
# GAP_BETWEEN = Inches(0.3)
# MARGIN = Inches(0.4)

# # Mapping from number to photo label
# mapping = {
#     1: "BUILDING PHOTO",
#     2: "CENTER MAIN GATE PHOTOGRAPH WITH COLLEGE_SCHOOL NAME & ADDRESS",
#     3: "AFFILIATION CERTIFICATE PHOTOGRAPH",
#     4: "CENTER HEAD BUSINESS CARD PHOTO OR COLLEGE_ SCHOOL ADDRESS PHOTO",
#     5: "SEATING PLAN PHOTO _ NODES PLAN PHOTO (LAB WISE)",
#     6: "LAB PHOTOS (CAPTURE IMAGE WITH DIFFERENT-DIFFERENT ANGLES)",
#     7: "FLOOR WISE LAB PHOTOS",
#     8: "LIFT PHOTO",
#     9: "RAMP PHOTO",
#     10: "DESK PHOTOS",
#     11: "CHAIR PHOTO",
#     12: "REGISTRATION DESK PHOTOS",
#     13: "CCTV CAMERA'S PHOTO",
#     14: "CCTV DISPLAY MONITOR PHOTO",
#     15: "CCTV DVR_NVR PHOTOGRAPH (OF MAKE & MODEL)",
#     16: "NETWORKING PLAN _ HUB ROOM PHOTO",
#     17: "SERVER ROOM PHOTOGRAPH",
#     18: "NETWORKING RACKS WITH COOLING",
#     19: "SYSTEM CONFIGURATION PHOTO",
#     20: "SYSTEM MONITOR PHOTO",
#     21: "INVOICE OR LICENSE PURCHASE (BILL OR AMC COPY REQUIRED)",
#     22: "MY COMPUTER PHOTO WITH DRIVE DETAILS",
#     23: "HDD CAPACITY IMAGE",
#     24: "PRINTER PHOTO",
#     25: "IT INVENTORY LIST",
#     26: "UPS PHOTOS WITH KVA DETAILS",
#     27: "BATTERY PHOTO WITH AH",
#     28: "DG PHOTOS ALONG WITH KVA DETAILS",
#     29: "POWER DIAGRAM (COPY REQUIRED)",
#     30: "SWITCHES PHOTO WITH MAKE AND MODEL",
#     31: "IP PHOTO",
#     32: "ROUTER PHOTO",
#     33: "SPEED TEST PHOTO",
#     34: "NETWORK DIAGRAM (COPY REQUIRED)",
#     35: "AIR CONDITION PHOTO (LAB & SERVER ROOM)",
#     36: "WATER COOLER OR RO PHOTOGRAPH",
#     37: "PARKING AREA PHOTO",
#     38: "TOILET PHOTO",
#     39: "TOILET PHOTO OF PH FRIENDLY",
#     40: "WHEELCHAIR PHOTO",
#     41: "BAGGAGE AREA PHOTO WITH PROPER RACKS FOR KEEPING BAGGAGES",
#     42: "FIREWALL PHOTO",
#     43: "NETWORK CABLING WITH WIRE TAGGING",
#     44: "CORE SWITCH CONNECTIVITY PHOTO",
#     45: "IP SERIES PHOTO",
#     46: "PING CHECK",
#     47: "ACCESS CONTROL - SERVER ROOM_ LAB",
#     48: "KEY BOARD PHOTO",
#     49: "MOUSE PHOTO",
#     50: "SCANNER PHOTO",
#     51: "BUFFER SYSTEM COUNT & PHOTO",
#     52: "SERVICE CERTIFICATE (DG & UPS) PHOTOS REQUIRED",
#     53: "SINGLE LINE DIAGRAM FOR POWER INFRA AND LOAD DISTRIBUTION",
#     54: "DG EXHAUST PIPE PHOTO",
#     55: "POWER OUTLET PHOTO",
#     56: "MEDICAL ROOM PHOTO",
#     57: "FIRE EXTINGUISHER PHOTO WITH EXPIRY CERTIFICATE",
#     58: "FIRE NOC REQUIRED",
#     59: "EMERGENCY EXIT PLAN PHOTO",
#     60: "SIGNAGE PHOTOS",
#     81: "Labs having Visible Seat Numbers"
# }
# for i in range(61, 81):
#     mapping[i] = f"other {i - 60}"

# # ======================
# # Core Functions
# # ======================
# def scale_image(img_path, max_width, max_height):
#     with Image.open(img_path) as img:
#         width, height = img.size
#         aspect = height / width
        
#         width_based = min(max_width, (max_height - GAP_BETWEEN) / aspect)
#         height_based = min(max_height - GAP_BETWEEN, max_width * aspect)
        
#         if (width_based * (width_based * aspect)) > (height_based * (height_based / aspect)):
#             return width_based, width_based * aspect
#         else:
#             return height_based / aspect, height_based

# def create_ppt(photo_paths, output_path):
#     prs = Presentation()
#     prs.slide_width = SLIDE_WIDTH
#     prs.slide_height = SLIDE_HEIGHT
    
#     for img_path in photo_paths:
#         slide = prs.slides.add_slide(prs.slide_layouts[6])
#         title = os.path.splitext(os.path.basename(img_path))[0]
        
#         textbox = slide.shapes.add_textbox(
#             left=Inches(0.5),
#             top=Inches(0.3),
#             width=SLIDE_WIDTH - Inches(1),
#             height=TITLE_HEIGHT
#         )
#         p = textbox.text_frame.add_paragraph()
#         p.text = title
#         p.alignment = PP_ALIGN.CENTER
#         p.font.name = "Arial"
#         p.font.size = Pt(24)
#         p.font.bold = True
        
#         max_width = SLIDE_WIDTH - (2 * MARGIN)
#         max_height = SLIDE_HEIGHT - TITLE_HEIGHT - GAP_BETWEEN
        
#         img_width, img_height = scale_image(img_path, max_width, max_height)
#         left = (SLIDE_WIDTH - img_width) / 2
#         top = TITLE_HEIGHT + GAP_BETWEEN
#         slide.shapes.add_picture(img_path, left, top, width=img_width, height=img_height)
    
#     prs.save(output_path)

# def create_pdf(photo_paths, output_path):
#     c = canvas.Canvas(output_path, pagesize=PDF_PAGE_SIZE)
#     page_width, page_height = PDF_PAGE_SIZE
#     title_font_size = 24
#     gap_points = GAP_BETWEEN.inches * 72
    
#     for img_path in photo_paths:
#         c.showPage()
#         title = os.path.splitext(os.path.basename(img_path))[0]
#         c.setFont("Helvetica-Bold", title_font_size)
#         title_y = page_height - 50
#         c.drawCentredString(page_width/2, title_y, title)
        
#         max_width = page_width - (MARGIN.inches * 72 * 2)
#         max_height = page_height - 70 - gap_points
        
#         with Image.open(img_path) as img:
#             img_width, img_height = scale_image(img_path, max_width, max_height)
        
#         x = (page_width - img_width) / 2
#         y = title_y - gap_points - img_height
#         c.drawImage(img_path, x, y, width=img_width, height=img_height, preserveAspectRatio=True)
    
#     c.save()

# # ======================
# # Streamlit UI
# # ======================
# st.set_page_config(page_title="Professional Photo Converter", layout="wide")
# st.title("📸 Smart Photo Converter")
# st.write("Upload photos to rename, generate a presentation and PDF.")

# uploaded_files = st.file_uploader(
#     "Upload Photos (numbered with IDs like 1, 2, ..., 81)",
#     type=["png", "jpg", "jpeg"],
#     accept_multiple_files=True
# )

# if uploaded_files:
#     temp_dir = tempfile.mkdtemp()
#     renamed_photo_paths = []
#     rename_log = []

#     for uploaded_file in uploaded_files:
#         original_name = uploaded_file.name
#         match = re.search(r"\b(\d{1,2}|81)\b", original_name)
        
#         if match:
#             number = int(match.group(1))
#             label = mapping.get(number, f"other {number}")
#             new_name = f"{number:02d} - {label}{os.path.splitext(original_name)[1]}"
#         else:
#             new_name = original_name  # No renaming if no number found
        
#         file_path = os.path.join(temp_dir, new_name)
#         with open(file_path, "wb") as f:
#             f.write(uploaded_file.getbuffer())
        
#         renamed_photo_paths.append(file_path)
#         rename_log.append([original_name, new_name])

#     # Preview
#     st.subheader("Photo Preview")
#     cols = st.columns(3)
#     for idx, path in enumerate(renamed_photo_paths):
#         try:
#             with warnings.catch_warnings():
#                 warnings.filterwarnings("ignore", category=DeprecationWarning)
#                 cols[idx % 3].image(path, use_column_width=True, caption=os.path.basename(path))
#         except Exception as e:
#             st.error(f"Error previewing {os.path.basename(path)}: {str(e)}")

#     # Output Options
#     st.subheader("Output Settings")
#     output_name = st.text_input("Base filename:", "MyPresentation")
#     output_format = st.radio("Format:", ["PPTX", "PDF", "Both"])

#     if st.button("🚀 Generate Files"):
#         with st.spinner("Generating files..."):
#             try:
#                 generated_files = []

#                 if output_format in ["PPTX", "Both"]:
#                     ppt_path = os.path.join(temp_dir, f"{output_name}.pptx")
#                     create_ppt(renamed_photo_paths, ppt_path)
#                     generated_files.append(ppt_path)

#                 if output_format in ["PDF", "Both"]:
#                     pdf_path = os.path.join(temp_dir, f"{output_name}.pdf")
#                     create_pdf(renamed_photo_paths, pdf_path)
#                     generated_files.append(pdf_path)

#                 log_csv_path = os.path.join(temp_dir, "rename_log.csv")
#                 with open(log_csv_path, mode="w", newline="", encoding="utf-8") as f:
#                     writer = csv.writer(f)
#                     writer.writerow(["Original Filename", "Renamed Filename"])
#                     writer.writerows(rename_log)
#                 generated_files.append(log_csv_path)

#                 st.success("✅ Files generated!")
#                 for file_path in generated_files:
#                     with open(file_path, "rb") as f:
#                         st.download_button(
#                             label=f"📥 Download {os.path.basename(file_path)}",
#                             data=f.read(),
#                             file_name=os.path.basename(file_path),
#                             mime="application/octet-stream"
#                         )
#             except Exception as e:
#                 st.error(f"⚠️ Error: {str(e)}")
#             finally:
#                 shutil.rmtree(temp_dir, ignore_errors=True)



















import streamlit as st
import tempfile
import os
import re
import csv
import shutil
import warnings
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from PIL import Image
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas

# ======================
# Configuration
# ======================
SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)
PDF_PAGE_SIZE = A4
TITLE_HEIGHT = Inches(0.8)
GAP_BETWEEN = Inches(0.3)
MARGIN = Inches(0.4)

# ======================
# Natural Sorting Utility
# ======================
def natural_sort_key(s):
    # Split string into list of ints and strings for natural sorting
    _nsre = re.compile(r'(\d+)')
    return [int(text) if text.isdigit() else text.lower() for text in _nsre.split(s)]

# ======================
# Photo Scaling for PPT/PDF
# ======================
def scale_image(img_path, max_width, max_height):
    with Image.open(img_path) as img:
        width, height = img.size
        aspect = height / width
        
        width_based = min(max_width, (max_height - GAP_BETWEEN) / aspect)
        height_based = min(max_height - GAP_BETWEEN, max_width * aspect)
        
        if (width_based * (width_based * aspect)) > (height_based * (height_based / aspect)):
            return width_based, width_based * aspect
        else:
            return height_based / aspect, height_based

# ======================
# PPT Creation
# ======================
def create_ppt(photo_paths, output_path):
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Sort photos naturally by filename
    photo_paths = sorted(photo_paths, key=lambda p: natural_sort_key(os.path.basename(p)))
    
    for img_path in photo_paths:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        title = os.path.splitext(os.path.basename(img_path))[0]
        textbox = slide.shapes.add_textbox(
            left=Inches(0.5),
            top=Inches(0.3),
            width=SLIDE_WIDTH - Inches(1),
            height=TITLE_HEIGHT
        )
        p = textbox.text_frame.add_paragraph()
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Arial"
        p.font.size = Pt(24)
        p.font.bold = True
        
        max_width = SLIDE_WIDTH - (2 * MARGIN)
        max_height = SLIDE_HEIGHT - TITLE_HEIGHT - GAP_BETWEEN
        
        img_width, img_height = scale_image(img_path, max_width, max_height)
        
        left = (SLIDE_WIDTH - img_width) / 2
        top = TITLE_HEIGHT + GAP_BETWEEN
        
        slide.shapes.add_picture(
            img_path,
            left, top,
            width=img_width,
            height=img_height
        )
    
    prs.save(output_path)

# ======================
# PDF Creation
# ======================
def create_pdf(photo_paths, output_path):
    c = canvas.Canvas(output_path, pagesize=PDF_PAGE_SIZE)
    page_width, page_height = PDF_PAGE_SIZE
    title_font_size = 24
    gap_points = GAP_BETWEEN.inches * 72

    photo_paths = sorted(photo_paths, key=lambda p: natural_sort_key(os.path.basename(p)))

    for img_path in photo_paths:
        c.showPage()
        
        title = os.path.splitext(os.path.basename(img_path))[0]
        c.setFont("Helvetica-Bold", title_font_size)
        title_y = page_height - 50
        c.drawCentredString(page_width/2, title_y, title)
        
        max_width = page_width - (MARGIN.inches * 72 * 2)
        max_height = page_height - 70 - gap_points
        
        with Image.open(img_path) as img:
            img_width, img_height = scale_image(img_path, max_width, max_height)
        
        x = (page_width - img_width) / 2
        y = title_y - gap_points - img_height
        
        c.drawImage(img_path, x, y, 
                   width=img_width, 
                   height=img_height,
                   preserveAspectRatio=True)
    
    c.save()

# ======================
# Renaming Logic
# ======================
mapping = {
    1: "BUILDING PHOTO",
    2: "CENTER MAIN GATE PHOTOGRAPH WITH COLLEGE_SCHOOL NAME & ADDRESS",
    3: "AFFILIATION CERTIFICATE PHOTOGRAPH",
    4: "CENTER HEAD BUSINESS CARD PHOTO OR COLLEGE_ SCHOOL ADDRESS PHOTO",
    5: "SEATING PLAN PHOTO _ NODES PLAN PHOTO (LAB WISE)",
    6: "LAB PHOTOS (CAPTURE IMAGE WITH DIFFERENT-DIFFERENT ANGLES)",
    7: "FLOOR WISE LAB PHOTOS",
    8: "LIFT PHOTO",
    9: "RAMP PHOTO",
    10: "DESK PHOTOS",
    11: "CHAIR PHOTO",
    12: "REGISTRATION DESK PHOTOS",
    13: "CCTV CAMERA'S PHOTO",
    14: "CCTV DISPLAY MONITOR PHOTO",
    15: "CCTV DVR_NVR PHOTOGRAPH (OF MAKE & MODEL)",
    16: "NETWORKING PLAN _ HUB ROOM PHOTO",
    17: "SERVER ROOM PHOTOGRAPH",
    18: "NETWORKING RACKS WITH COOLING",
    19: "SYSTEM CONFIGURATION PHOTO",
    20: "SYSTEM MONITOR PHOTO",
    21: "INVOICE OR LICENSE PURCHASE (BILL OR AMC COPY REQUIRED)",
    22: "MY COMPUTER PHOTO WITH DRIVE DETAILS",
    23: "HDD CAPACITY IMAGE",
    24: "PRINTER PHOTO",
    25: "IT INVENTORY LIST",
    26: "UPS PHOTOS WITH KVA DETAILS",
    27: "BATTERY PHOTO WITH AH",
    28: "DG PHOTOS ALONG WITH KVA DETAILS",
    29: "POWER DIAGRAM (COPY REQUIRED)",
    30: "SWITCHES PHOTO WITH MAKE AND MODEL",
    31: "IP PHOTO",
    32: "ROUTER PHOTO",
    33: "SPEED TEST PHOTO",
    34: "NETWORK DIAGRAM (COPY REQUIRED)",
    35: "AIR CONDITION PHOTO (LAB & SERVER ROOM)",
    36: "WATER COOLER OR RO PHOTOGRAPH",
    37: "PARKING AREA PHOTO",
    38: "TOILET PHOTO",
    39: "TOILET PHOTO OF PH FRIENDLY",
    40: "WHEELCHAIR PHOTO",
    41: "BAGGAGE AREA PHOTO WITH PROPER RACKS FOR KEEPING BAGGAGES",
    42: "FIREWALL PHOTO",
    43: "NETWORK CABLING WITH WIRE TAGGING",
    44: "CORE SWITCH CONNECTIVITY PHOTO",
    45: "IP SERIES PHOTO",
    46: "PING CHECK",
    47: "ACCESS CONTROL - SERVER ROOM_ LAB",
    48: "KEY BOARD PHOTO",
    49: "MOUSE PHOTO",
    50: "SCANNER PHOTO",
    51: "BUFFER SYSTEM COUNT & PHOTO",
    52: "SERVICE CERTIFICATE (DG & UPS) PHOTOS REQUIRED",
    53: "SINGLE LINE DIAGRAM FOR POWER INFRA AND LOAD DISTRIBUTION",
    54: "DG EXHAUST PIPE PHOTO",
    55: "POWER OUTLET PHOTO",
    56: "MEDICAL ROOM PHOTO",
    57: "FIRE EXTINGUISHER PHOTO WITH EXPIRY CERTIFICATE",
    58: "FIRE NOC REQUIRED",
    59: "EMERGENCY EXIT PLAN PHOTO",
    60: "SIGNAGE PHOTOS"
}

# Add "other X" entries for 61–80
for i in range(61, 81):
    mapping[i] = f"other {i - 60}"

# Special case
mapping[81] = "Labs having Visible Seat Numbers"

image_extensions = ('.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp')

def rename_photos_in_folder(folder_path, dry_run=False):
    rename_data = []
    status_logs = []
    total_files = 0
    renamed_count = 0

    for dirpath, dirnames, filenames in os.walk(folder_path):
        if os.path.basename(dirpath).lower() == 'photo':
            status_logs.append(f"📁 Processing Photo folder: {dirpath}")
            for filename in filenames:
                if not filename.lower().endswith(image_extensions):
                    continue

                total_files += 1

                file_path = os.path.join(dirpath, filename)
                base_name, ext = os.path.splitext(filename)

                match = re.match(r'^(\d+)', base_name)
                if not match:
                    status_logs.append(f"⚠️ No number prefix in file: {filename} - Skipping")
                    continue
                
                num = int(match.group(1))
                rest = base_name[len(match.group(1)):].lstrip('_ -')

                new_base = mapping.get(num, f"other {num}")  # default if not found
                new_name = new_base
                if rest:
                    new_name += " " + rest
                new_name += ext

                new_path = os.path.join(dirpath, new_name)

                if dry_run:
                    status_logs.append(f"Dry run: Would rename {filename} → {new_name}")
                else:
                    try:
                        os.rename(file_path, new_path)
                        status_logs.append(f"Renamed: {filename} → {new_name}")
                        renamed_count += 1
                        rename_data.append([filename, new_name, dirpath])
                    except Exception as e:
                        status_logs.append(f"❌ Error renaming {filename}: {e}")

    summary = f"Total files found: {total_files}, Renamed: {renamed_count}"
    return rename_data, status_logs, summary

# ======================
# Streamlit UI
# ======================
def main():
    st.set_page_config(page_title="Photo Renamer & PPT Creator", layout="wide")

    st.title("Photo Renamer & PPT Creator")

    menu = ["Rename Photos", "Create PPT from Photos"]
    choice = st.sidebar.selectbox("Choose an action", menu)

    if choice == "Rename Photos":
        st.header("Rename Photos in 'Photo' Folder")

        folder_path = st.text_input("Enter base folder path containing 'Photo' subfolders:", "")

        dry_run = st.checkbox("Dry Run (show what will be renamed, no changes)", True)

        if st.button("Start Renaming") and folder_path:
            if not os.path.exists(folder_path):
                st.error("Folder path does not exist!")
            else:
                with st.spinner("Renaming photos..."):
                    rename_data, logs, summary = rename_photos_in_folder(folder_path, dry_run=dry_run)
                st.success(summary)
                st.text_area("Logs:", value="\n".join(logs), height=300)

                if rename_data:
                    st.markdown("### Renamed Files Summary")
                    st.dataframe(rename_data)

    elif choice == "Create PPT from Photos":
        st.header("Create PowerPoint from Photos")

        uploaded_files = st.file_uploader("Upload photo files", accept_multiple_files=True, type=['jpg','jpeg','png','gif','bmp','tiff','webp'])

        ppt_name = st.text_input("Enter PPT filename (without extension):", "output_presentation")

        if uploaded_files and ppt_name:
            temp_dir = tempfile.mkdtemp()

            photo_paths = []
            for file in uploaded_files:
                temp_file_path = os.path.join(temp_dir, file.name)
                with open(temp_file_path, "wb") as f:
                    f.write(file.getbuffer())
                photo_paths.append(temp_file_path)

            ppt_path = os.path.join(temp_dir, ppt_name + ".pptx")

            with st.spinner("Creating PPT..."):
                create_ppt(photo_paths, ppt_path)
            
            with open(ppt_path, "rb") as f:
                ppt_bytes = f.read()

            st.download_button("Download PPTX", data=ppt_bytes, file_name=ppt_name + ".pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

            shutil.rmtree(temp_dir)

if __name__ == "__main__":
    main()

